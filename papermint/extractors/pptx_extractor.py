"""PowerPoint extractor module for PaperMint."""

from __future__ import annotations

import io
import logging
from collections.abc import Iterator
from typing import Any

from pptx import Presentation

from papermint.errors import CorruptedDocumentError
from papermint.extractors.base import BaseExtractor, ExtractedDocument

logger = logging.getLogger(__name__)

#: MSO shape type identifier for a grouped shape.
_GROUP_SHAPE_TYPE = 6


class PPTXExtractor(BaseExtractor):
    """Extractor for PowerPoint (.pptx) files using python-pptx."""

    name = "PowerPoint"
    mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )
    extensions = frozenset({"pptx"})

    def _walk_shapes(self, shapes: Any) -> Iterator[Any]:
        """Yield every shape, descending into grouped shapes.

        Args:
            shapes: A python-pptx shape collection.

        Yields:
            Each leaf shape in the tree.
        """
        for shape in shapes:
            if getattr(shape, "shape_type", None) == _GROUP_SHAPE_TYPE:
                yield from self._walk_shapes(shape.shapes)
            else:
                yield shape

    def _slide_text(self, slide: Any) -> str:
        """Collect all text belonging to a single slide.

        Args:
            slide: A python-pptx slide.

        Returns:
            The slide's text, including tables and speaker notes.
        """
        lines: list[str] = []

        for shape in self._walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        lines.append(" | ".join(dict.fromkeys(cells)))

        if getattr(slide, "has_notes_slide", False):
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame else ""
            if notes:
                lines.append(notes)

        return "\n".join(lines)

    def extract_text(self, file_bytes: bytes) -> str:
        """Extract text from a PowerPoint presentation.

        Args:
            file_bytes: The raw .pptx bytes.

        Returns:
            The presentation text, one slide after another.

        Raises:
            CorruptedDocumentError: If the archive cannot be opened.
        """
        return self.extract_document(file_bytes).text

    def extract_document(self, file_bytes: bytes) -> ExtractedDocument:
        """Extract text along with the slide count.

        Args:
            file_bytes: The raw .pptx bytes.

        Returns:
            An :class:`ExtractedDocument` including per-slide text.

        Raises:
            CorruptedDocumentError: If the archive cannot be opened.
        """
        try:
            presentation = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            logger.exception("Failed to open PPTX stream")
            raise CorruptedDocumentError(
                "This presentation could not be opened. It may be corrupted, or it "
                "may be an older .ppt file rather than .pptx.",
                remedy="Open it in PowerPoint and save it again as .pptx.",
            ) from exc

        slides = tuple(self._slide_text(slide) for slide in presentation.slides)
        return ExtractedDocument(
            text="\n\n".join(s for s in slides if s),
            page_count=len(slides),
            pages=slides,
        )
